#!/usr/bin/env python3
"""AI-ROI-simulator_v1 紹介・使い方説明 PowerPoint 生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# ─── ブランドカラー定義 ───
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)        # 濃紺背景
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)    # ブルーアクセント
ACCENT_CYAN = RGBColor(0x06, 0xB6, 0xD4)    # シアン
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)   # グリーン
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)   # アンバー
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)
CARD_BG = RGBColor(0x1E, 0x29, 0x3B)        # カード背景
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)       # 明るい背景（印刷向け）

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── ヘルパー関数 ───
def add_dark_bg(slide):
    """スライド全体にダーク背景を設定"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_light_bg(slide):
    """スライド全体にライト背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_BG


def add_gradient_bar(slide, left, top, width, height, color1, color2):
    """グラデーション風のアクセントバーを追加"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color1
    shape.line.fill.background()


def add_rounded_card(slide, left, top, width, height, fill_color=CARD_BG):
    """角丸カードを追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    """テキストボックスを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=16,
                       color=WHITE, line_spacing=1.5, font_name="Meiryo"):
    """複数行テキストを追加"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, txt_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = txt_color if txt_color else color
        p.font.bold = is_bold
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.4)
    return txBox


def add_page_number(slide, num, total):
    """ページ番号"""
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                 f"{num}/{total}", font_size=10, color=MID_GRAY,
                 alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12


# ============================================================
# SLIDE 1: タイトルスライド
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_dark_bg(slide)

# アクセントライン上部
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

# メインタイトル
add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
             "AI-ROI Simulator", font_size=54, color=WHITE, bold=True)

# サブタイトル
add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
             "AI導入ROI試算ツール — 機能説明 & 使い方ガイド",
             font_size=28, color=ACCENT_CYAN)

# 説明
add_text_box(slide, Inches(1), Inches(4.2), Inches(9), Inches(1.0),
             "初回商談でクライアントに「御社のAI導入ROI」を即座にシミュレーション。\n提案の説得力を劇的に高める、AIコンサルタント専用ツール。",
             font_size=16, color=LIGHT_GRAY)

# ブランドバッジ
card = add_rounded_card(slide, Inches(1), Inches(5.8), Inches(3.2), Inches(0.6), ACCENT_BLUE)
add_text_box(slide, Inches(1.2), Inches(5.85), Inches(3), Inches(0.5),
             "AI Transformation Architect", font_size=14, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

# バージョン
add_text_box(slide, Inches(9), Inches(6.5), Inches(3), Inches(0.4),
             "Version 1.0  |  2026", font_size=12, color=MID_GRAY,
             alignment=PP_ALIGN.RIGHT)

add_page_number(slide, 1, TOTAL_SLIDES)


# ============================================================
# SLIDE 2: ツール概要 — このツールとは？
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "OVERVIEW", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "AI-ROI Simulator とは？", font_size=36, color=WHITE, bold=True)

# 3カラム構成
cols = [
    ("商談クロージングツール",
     "初回商談30分で「御社のAI導入ROI」を\nリアルタイムにシミュレーション。\n数字の力で次回アポ・受注を引き出します。",
     ACCENT_BLUE),
    ("ポートフォリオ",
     "React + TypeScript で構築された\nプロフェッショナルなWebアプリ。\nAI変革アーキテクトとしての技術力の証明。",
     ACCENT_CYAN),
    ("差別化の武器",
     "日本市場初のインタラクティブ\nAI ROIシミュレーター。\n競合コンサルタントとの決定的な差別化要因。",
     ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(cols):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_card(slide, x, Inches(2.0), Inches(3.7), Inches(3.8))
    # アクセントライン
    add_gradient_bar(slide, x, Inches(2.0), Inches(3.7), Inches(0.05), color, color)
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(3.1), Inches(0.5),
                 title, font_size=20, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(3.0), Inches(3.1), Inches(2.5),
                 desc, font_size=14, color=LIGHT_GRAY)

# ボトムメッセージ
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             "「この数字の根拠は？」に即座に答えられる —— 前提条件の完全な透明性を確保",
             font_size=15, color=ACCENT_AMBER, bold=True)

add_page_number(slide, 2, TOTAL_SLIDES)


# ============================================================
# SLIDE 3: 主要機能一覧
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "FEATURES", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "主要機能一覧（MVP）", font_size=36, color=WHITE, bold=True)

features = [
    ("F1", "パラメータ入力UI", "業種・規模・対象業務をステッパー形式で入力。\n商談の流れに沿ったガイド付きフォーム。", ACCENT_BLUE),
    ("F2", "ROIダッシュボード", "ヒーローナンバー + KPIカード + グラフで\n試算結果を一目で把握。30秒で理解可能。", ACCENT_CYAN),
    ("F3", "3シナリオ比較", "保守的 / 標準 / 楽観的の3シナリオを\n並列表示。信頼性の担保と説得力の両立。", ACCENT_GREEN),
    ("F4", "業種別テンプレート", "10業種のプリセットパラメータ。\n業種選択だけで即座に試算を開始。", ACCENT_AMBER),
    ("F9", "感度分析", "スライダー操作でROIがリアルタイム変動。\nExcelでは不可能なインタラクティブ体験。", ACCENT_BLUE),
    ("F11", "導入あり/なし比較", "AI導入前後の3〜5年推移グラフ。\n経営層が最も直感的に理解できる表示。", ACCENT_CYAN),
]

for i, (fid, title, desc, color) in enumerate(features):
    row = i // 3
    col = i % 3
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.9 + row * 2.6)

    card = add_rounded_card(slide, x, y, Inches(3.7), Inches(2.2))

    # 機能ID バッジ
    badge = add_rounded_card(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.35), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.22), Inches(0.6), Inches(0.3),
                 fid, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.0), y + Inches(0.2), Inches(2.5), Inches(0.4),
                 title, font_size=17, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.75), Inches(3.1), Inches(1.3),
                 desc, font_size=12, color=LIGHT_GRAY)

add_page_number(slide, 3, TOTAL_SLIDES)


# ============================================================
# SLIDE 4: 使い方 STEP 1 — パラメータ入力
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "HOW TO USE — STEP 1", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "Step 1: クライアント情報を入力", font_size=36, color=WHITE, bold=True)

# 左側：入力フォームのモックアップ
card = add_rounded_card(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(5.0))

form_items = [
    ("業種を選択", "製造業 / 小売・EC / 金融 / IT / 医療 / 物流 / 不動産 / 教育 / メディア / 専門サービス"),
    ("従業員数", "〜50名 / 51-200名 / 201-500名 / 501-1,000名 / 1,001-5,000名 / 5,001名〜"),
    ("年間売上規模", "〜1億円 / 1-10億円 / 10-50億円 / 50-100億円 / 100-500億円 / 500億円〜"),
    ("AI導入対象業務", "カスタマーサポート / 営業・マーケ / データ分析 / 文書管理 / 経理 / 人事 / 品質管理 / SCM / 広告運用 / R&D"),
]

for i, (label, options) in enumerate(form_items):
    y = Inches(2.2 + i * 1.15)
    add_text_box(slide, Inches(1.2), y, Inches(4.8), Inches(0.35),
                 label, font_size=15, color=ACCENT_CYAN, bold=True)
    # 入力フィールド風
    field = add_rounded_card(slide, Inches(1.2), y + Inches(0.35), Inches(4.6), Inches(0.55),
                             RGBColor(0x15, 0x20, 0x33))
    add_text_box(slide, Inches(1.4), y + Inches(0.38), Inches(4.2), Inches(0.45),
                 options, font_size=9, color=MID_GRAY)

# 右側：ポイント解説
add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.5),
             "入力のポイント", font_size=22, color=WHITE, bold=True)

points = [
    ("二段構え設計", "必須入力（4項目）だけで概算が出ます。\n任意入力でさらに精度を上げることも可能。", ACCENT_BLUE),
    ("商談で聞ける範囲", "必須項目はすべて初回商談のヒアリングで\n自然に聞ける情報のみに限定。", ACCENT_GREEN),
    ("業種テンプレート", "業種を選ぶだけで典型的なパラメータが\n自動プリセット。入力の手間を最小化。", ACCENT_AMBER),
    ("リアルタイム反映", "入力と同時にプレビューが更新。\n商談中のテンポを損なわない設計。", ACCENT_CYAN),
]

for i, (title, desc, color) in enumerate(points):
    y = Inches(2.7 + i * 1.15)
    add_gradient_bar(slide, Inches(7.0), y, Inches(0.05), Inches(0.8), color, color)
    add_text_box(slide, Inches(7.3), y, Inches(5), Inches(0.35),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(7.3), y + Inches(0.35), Inches(5), Inches(0.5),
                 desc, font_size=11, color=LIGHT_GRAY)

add_page_number(slide, 4, TOTAL_SLIDES)


# ============================================================
# SLIDE 5: 使い方 STEP 2 — ROIダッシュボード
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "HOW TO USE — STEP 2", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "Step 2: ROIダッシュボードで結果を確認", font_size=36, color=WHITE, bold=True)

# ヒーローナンバー
hero_card = add_rounded_card(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.5),
                              RGBColor(0x1A, 0x25, 0x3D))
add_gradient_bar(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(0.05), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(1.2), Inches(2.0), Inches(4), Inches(0.4),
             "御社のAI導入による年間効果（標準シナリオ）", font_size=14, color=MID_GRAY)
add_text_box(slide, Inches(1.2), Inches(2.4), Inches(5), Inches(0.8),
             "¥ 2,400 万円", font_size=44, color=ACCENT_GREEN, bold=True)
add_text_box(slide, Inches(1.2), Inches(3.0), Inches(5), Inches(0.3),
             "年間コスト削減額", font_size=13, color=LIGHT_GRAY)

# KPIカード群
kpis = [
    ("3年ROI", "320%", ACCENT_BLUE),
    ("投資回収", "8.5 ヶ月", ACCENT_CYAN),
    ("売上増加", "¥1,800万/年", ACCENT_GREEN),
    ("FTE換算", "3.2 人分", ACCENT_AMBER),
]

for i, (label, value, color) in enumerate(kpis):
    x = Inches(0.8 + i * 3.05)
    y = Inches(3.7)
    kpi_card = add_rounded_card(slide, x, y, Inches(2.7), Inches(1.3))
    add_gradient_bar(slide, x, y, Inches(2.7), Inches(0.04), color, color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.3), Inches(0.3),
                 label, font_size=12, color=MID_GRAY)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.5), Inches(2.3), Inches(0.6),
                 value, font_size=28, color=color, bold=True)

# 説明テキスト
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(0.4),
             "「30秒ルール」— 経営層が画面を見て30秒で「うちにとっての効果」を理解できるUI設計",
             font_size=14, color=ACCENT_AMBER, bold=True)

desc_text = ("最上部のヒーローナンバーで最も重要な数字を大きく表示。KPIカードで主要指標を一覧。"
             "すべての数字はクライアントが入力した情報に基づくパーソナライズされた試算です。")
add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.8),
             desc_text, font_size=12, color=LIGHT_GRAY)

# 免責事項
add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
             "※ 試算値は公開調査レポートおよび類似プロジェクト実績に基づく推定値です。前提条件はダッシュボード下部に明示されます。",
             font_size=10, color=MID_GRAY)

add_page_number(slide, 5, TOTAL_SLIDES)


# ============================================================
# SLIDE 6: 使い方 STEP 3 — シナリオ比較
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "HOW TO USE — STEP 3", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "Step 3: 3シナリオで信頼性を担保", font_size=36, color=WHITE, bold=True)

# 3つのシナリオカード
scenarios = [
    ("保守的", "Conservative",
     [("年間効果", "¥1,200万"), ("3年ROI", "180%"), ("回収期間", "14ヶ月")],
     MID_GRAY, "最低限期待できる効果。\nすべてのパラメータに下限値を適用。"),
    ("標準", "Standard",
     [("年間効果", "¥2,400万"), ("3年ROI", "320%"), ("回収期間", "8.5ヶ月")],
     ACCENT_BLUE, "一般的に期待できる効果。\n業界平均のベンチマーク値を適用。"),
    ("楽観的", "Optimistic",
     [("年間効果", "¥4,100万"), ("3年ROI", "580%"), ("回収期間", "5ヶ月")],
     ACCENT_GREEN, "最大限の効果が出た場合。\nすべてのパラメータに上限値を適用。"),
]

for i, (name, eng, kpis_list, color, desc) in enumerate(scenarios):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_card(slide, x, Inches(1.9), Inches(3.7), Inches(4.5))
    add_gradient_bar(slide, x, Inches(1.9), Inches(3.7), Inches(0.05), color, color)

    add_text_box(slide, x + Inches(0.3), Inches(2.15), Inches(3.1), Inches(0.5),
                 name, font_size=24, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.6), Inches(3.1), Inches(0.3),
                 eng, font_size=11, color=MID_GRAY)

    for j, (kpi_label, kpi_val) in enumerate(kpis_list):
        y = Inches(3.15 + j * 0.75)
        add_text_box(slide, x + Inches(0.3), y, Inches(1.8), Inches(0.3),
                     kpi_label, font_size=12, color=MID_GRAY)
        add_text_box(slide, x + Inches(2.0), y, Inches(1.5), Inches(0.3),
                     kpi_val, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)

    add_text_box(slide, x + Inches(0.3), Inches(5.4), Inches(3.1), Inches(0.8),
                 desc, font_size=11, color=LIGHT_GRAY)

# ボトムメッセージ
add_text_box(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
             "「保守的に見積もっても年間1,200万円の効果が見込めます」— この一言がクライアントの信頼を勝ち取る",
             font_size=14, color=ACCENT_AMBER, bold=True)

add_page_number(slide, 6, TOTAL_SLIDES)


# ============================================================
# SLIDE 7: 使い方 STEP 4 — 感度分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "HOW TO USE — STEP 4", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "Step 4: 感度分析でインタラクティブに探索", font_size=36, color=WHITE, bold=True)

# 左側：スライダーUI モックアップ
card = add_rounded_card(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.8))

sliders = [
    ("AI自動化率", "25%", "15% ← → 40%", ACCENT_BLUE),
    ("導入投資額", "¥2,000万", "¥1,000万 ← → ¥3,000万", ACCENT_CYAN),
    ("効果発現期間", "6ヶ月", "3ヶ月 ← → 18ヶ月", ACCENT_GREEN),
    ("年間運用費率", "25%", "15% ← → 35%", ACCENT_AMBER),
]

for i, (label, value, slider_range, color) in enumerate(sliders):
    y = Inches(2.3 + i * 1.1)
    add_text_box(slide, Inches(1.2), y, Inches(2.5), Inches(0.3),
                 label, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.8), y, Inches(2), Inches(0.3),
                 value, font_size=13, color=color, bold=True, alignment=PP_ALIGN.RIGHT)
    # スライダーバー（背景）
    bar_bg = add_rounded_card(slide, Inches(1.2), y + Inches(0.35), Inches(4.6), Inches(0.12),
                               RGBColor(0x15, 0x20, 0x33))
    # スライダーバー（値）
    bar_val = add_rounded_card(slide, Inches(1.2), y + Inches(0.35), Inches(2.3), Inches(0.12),
                                color)
    # レンジ表示
    add_text_box(slide, Inches(1.2), y + Inches(0.5), Inches(4.6), Inches(0.25),
                 slider_range, font_size=9, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# 右側：リアルタイム更新の説明
add_text_box(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.5),
             "Excelとの決定的な違い", font_size=22, color=WHITE, bold=True)

right_points = [
    "スライダーを動かすと、グラフとKPIがリアルタイムに更新",
    "「もし自動化率がもう少し低かったら？」\n→ その場でスライダーを動かして即座に回答",
    "商談中の「What-if」質問に瞬時に対応可能",
    "各パラメータがROIに与える影響度を\nトルネードチャートで可視化",
    "Excel操作不要。クリック＆ドラッグだけで探索完了",
]

for i, point in enumerate(right_points):
    y = Inches(2.7 + i * 0.85)
    add_text_box(slide, Inches(7.0), y, Inches(0.3), Inches(0.3),
                 "▸", font_size=14, color=ACCENT_CYAN)
    add_text_box(slide, Inches(7.4), y, Inches(5), Inches(0.7),
                 point, font_size=12, color=LIGHT_GRAY)

add_page_number(slide, 7, TOTAL_SLIDES)


# ============================================================
# SLIDE 8: 使い方 STEP 5 — 導入あり/なし比較グラフ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "HOW TO USE — STEP 5", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "Step 5: 3〜5年の将来推移を比較", font_size=36, color=WHITE, bold=True)

# チャートを追加
chart_data = CategoryChartData()
chart_data.categories = ['1年目', '2年目', '3年目', '4年目', '5年目']
chart_data.add_series('AI導入なし', (100, 102, 104, 106, 108))
chart_data.add_series('AI導入あり（標準）', (85, 78, 72, 68, 65))
chart_data.add_series('累計コスト効果', (15, 39, 71, 109, 152))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS, Inches(0.8), Inches(1.9),
    Inches(7.5), Inches(4.5), chart_data
)
chart = chart_frame.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(11)
chart.legend.font.color.rgb = LIGHT_GRAY

plot = chart.plots[0]
plot.has_data_labels = False

# 系列のスタイリング
series_colors = [MID_GRAY, ACCENT_BLUE, ACCENT_GREEN]
for i, series in enumerate(plot.series):
    series.format.line.color.rgb = series_colors[i]
    series.format.line.width = Pt(3)
    series.smooth = True

# チャート背景
chart.chart_style = 2

# 右側の説明
add_text_box(slide, Inches(8.8), Inches(2.0), Inches(4), Inches(0.4),
             "グラフの読み方", font_size=20, color=WHITE, bold=True)

graph_points = [
    ("グレーの線", "AI導入なしの場合の\nコスト推移（現状維持）", MID_GRAY),
    ("ブルーの線", "AI導入ありの場合の\nコスト推移（削減効果）", ACCENT_BLUE),
    ("グリーンの線", "累計コスト効果\n（差額の積み上げ）", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(graph_points):
    y = Inches(2.7 + i * 1.2)
    add_gradient_bar(slide, Inches(8.8), y, Inches(0.4), Inches(0.04), color, color)
    add_text_box(slide, Inches(8.8), y + Inches(0.1), Inches(3.8), Inches(0.3),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(8.8), y + Inches(0.4), Inches(3.8), Inches(0.5),
                 desc, font_size=11, color=LIGHT_GRAY)

add_text_box(slide, Inches(8.8), Inches(5.8), Inches(4), Inches(0.8),
             "経営層が最も直感的に理解できる\n「Before / After」の可視化",
             font_size=14, color=ACCENT_AMBER, bold=True)

add_page_number(slide, 8, TOTAL_SLIDES)


# ============================================================
# SLIDE 9: ROI計算モデルの概要
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "ROI MODEL", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "ROI計算モデルの仕組み", font_size=36, color=WHITE, bold=True)

# 3つのモデルカード
models = [
    ("Model A", "コスト削減型",
     "対象業務コスト × AI自動化率\n× 業種補正係数",
     "バックオフィス業務の自動化、\n問い合わせ対応の自動化",
     ACCENT_BLUE),
    ("Model B", "売上増加型",
     "現在売上 × 対象業務売上貢献率\n× AI効果係数",
     "営業最適化、広告最適化、\n需要予測",
     ACCENT_GREEN),
    ("Model C", "生産性向上型",
     "削減時間 × FTE換算\n× 平均人件費",
     "すべての業務に適用可能な\n汎用モデル（FTE換算）",
     ACCENT_CYAN),
]

for i, (mid, name, formula, use_case, color) in enumerate(models):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_card(slide, x, Inches(1.9), Inches(3.7), Inches(3.2))
    add_gradient_bar(slide, x, Inches(1.9), Inches(3.7), Inches(0.05), color, color)

    add_text_box(slide, x + Inches(0.3), Inches(2.15), Inches(3.1), Inches(0.35),
                 f"{mid}: {name}", font_size=18, color=color, bold=True)

    # 計算式（コードブロック風）
    formula_bg = add_rounded_card(slide, x + Inches(0.2), Inches(2.7), Inches(3.3), Inches(0.9),
                                   RGBColor(0x0A, 0x10, 0x20))
    add_text_box(slide, x + Inches(0.35), Inches(2.75), Inches(3.0), Inches(0.8),
                 formula, font_size=11, color=ACCENT_CYAN, font_name="Consolas")

    add_text_box(slide, x + Inches(0.3), Inches(3.8), Inches(3.1), Inches(0.35),
                 "適用場面:", font_size=11, color=MID_GRAY, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(4.1), Inches(3.1), Inches(0.8),
                 use_case, font_size=11, color=LIGHT_GRAY)

# 下部：統合計算の説明
integrated_card = add_rounded_card(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.5))
add_text_box(slide, Inches(1.2), Inches(5.6), Inches(4), Inches(0.4),
             "複合ROI計算", font_size=18, color=WHITE, bold=True)

add_text_box(slide, Inches(1.2), Inches(6.0), Inches(10.5), Inches(0.3),
             "総合年間効果 = Σ（各対象業務の [コスト削減額 + 売上増加額 + 生産性向上額]）",
             font_size=12, color=ACCENT_CYAN, font_name="Consolas")
add_text_box(slide, Inches(1.2), Inches(6.35), Inches(10.5), Inches(0.3),
             "3年ROI = (総合年間効果 × 3 - AI導入投資合計) / AI導入投資合計 × 100",
             font_size=12, color=ACCENT_CYAN, font_name="Consolas")

add_page_number(slide, 9, TOTAL_SLIDES)


# ============================================================
# SLIDE 10: データの信頼性・根拠
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "TRUST & TRANSPARENCY", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "データの信頼性と根拠", font_size=36, color=WHITE, bold=True)

# データソーステーブル
sources = [
    ("McKinsey Global AI Survey", "業種別AI導入効果の平均値", "★★★★★"),
    ("Accenture AI Profit Impact", "業種別利益率改善の予測", "★★★★☆"),
    ("MIT Sloan Management Review", "企業規模別AI導入ROI", "★★★★☆"),
    ("IDC / Gartner Reports", "IT投資規模のベンチマーク", "★★★★★"),
    ("Deloitte AI Institute", "業種別AI活用事例・効果", "★★★★☆"),
    ("自社5社AI導入実績", "特定領域のキャリブレーション", "★★★☆☆"),
]

# テーブルヘッダー
header_bg = add_rounded_card(slide, Inches(0.8), Inches(1.9), Inches(7), Inches(0.45),
                              ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.95), Inches(2.8), Inches(0.35),
             "データソース", font_size=12, color=WHITE, bold=True)
add_text_box(slide, Inches(3.8), Inches(1.95), Inches(2.5), Inches(0.35),
             "用途", font_size=12, color=WHITE, bold=True)
add_text_box(slide, Inches(6.5), Inches(1.95), Inches(1.2), Inches(0.35),
             "信頼度", font_size=12, color=WHITE, bold=True)

for i, (source, usage, rating) in enumerate(sources):
    y = Inches(2.4 + i * 0.5)
    bg_color = CARD_BG if i % 2 == 0 else RGBColor(0x15, 0x20, 0x33)
    row_bg = add_rounded_card(slide, Inches(0.8), y, Inches(7), Inches(0.45), bg_color)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(2.8), Inches(0.35),
                 source, font_size=11, color=WHITE)
    add_text_box(slide, Inches(3.8), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 usage, font_size=11, color=LIGHT_GRAY)
    add_text_box(slide, Inches(6.5), y + Inches(0.05), Inches(1.2), Inches(0.35),
                 rating, font_size=11, color=ACCENT_AMBER)

# 右側：信頼性担保の設計
add_text_box(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(0.5),
             "信頼性を担保する設計", font_size=20, color=WHITE, bold=True)

trust_points = [
    ("3シナリオ比較", "楽観的な数字だけでなく\n保守的シナリオを必ず提示", ACCENT_BLUE),
    ("前提条件の明示", "すべての試算結果に\n計算の前提条件を表示", ACCENT_GREEN),
    ("データソース開示", "ベンチマーク値の出典を\nいつでも確認可能", ACCENT_CYAN),
    ("免責事項の表示", "推定値であることを明記。\nコンサルとしての誠実さ", ACCENT_AMBER),
]

for i, (title, desc, color) in enumerate(trust_points):
    y = Inches(2.7 + i * 1.1)
    add_gradient_bar(slide, Inches(8.5), y, Inches(0.05), Inches(0.8), color, color)
    add_text_box(slide, Inches(8.8), y, Inches(3.8), Inches(0.3),
                 title, font_size=13, color=color, bold=True)
    add_text_box(slide, Inches(8.8), y + Inches(0.3), Inches(3.8), Inches(0.5),
                 desc, font_size=11, color=LIGHT_GRAY)

add_page_number(slide, 10, TOTAL_SLIDES)


# ============================================================
# SLIDE 11: 商談での活用フロー
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "MEETING FLOW", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "商談30分での活用フロー", font_size=36, color=WHITE, bold=True)

# タイムラインフロー
phases = [
    ("0-5分", "ヒアリング", "自己紹介\n課題のヒアリング\nツールは未使用",
     MID_GRAY, "通常の商談と同じ流れで\n自然にスタート"),
    ("5-15分", "業種・課題特定", "業種を選択\nAI適用業務を特定\nテンプレート選択",
     ACCENT_BLUE, "「御社の業種では、この業務に\nAIが最も効果的です」"),
    ("15-25分", "ROI提示 ★", "パラメータ入力\nリアルタイム試算\n感度分析",
     ACCENT_GREEN, "「御社の場合、年間2,400万円の\n効果が見込めます」"),
    ("25-30分", "クロージング", "シナリオ比較\n将来推移グラフ\n次回アポ取得",
     ACCENT_AMBER, "「詳細なレポートを\nお送りいたします」"),
]

for i, (time, title, actions, color, script) in enumerate(phases):
    x = Inches(0.5 + i * 3.2)
    y = Inches(2.0)

    # タイム表示
    time_badge = add_rounded_card(slide, x + Inches(0.5), y, Inches(1.5), Inches(0.4), color)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.03), Inches(1.5), Inches(0.35),
                 time, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # フェーズタイトル
    add_text_box(slide, x + Inches(0.1), y + Inches(0.6), Inches(2.8), Inches(0.4),
                 title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # アクションカード
    action_card = add_rounded_card(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.8), Inches(1.5))
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(2.4), Inches(1.3),
                 actions, font_size=11, color=LIGHT_GRAY)

    # トークスクリプト
    script_card = add_rounded_card(slide, x + Inches(0.1), y + Inches(2.8), Inches(2.8), Inches(1.2),
                                    RGBColor(0x15, 0x20, 0x33))
    add_text_box(slide, x + Inches(0.2), y + Inches(2.85), Inches(2.6), Inches(0.2),
                 "トークスクリプト:", font_size=9, color=MID_GRAY, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(3.1), Inches(2.6), Inches(0.8),
                 script, font_size=10, color=ACCENT_CYAN)

    # 矢印（最後以外）
    if i < 3:
        add_text_box(slide, x + Inches(3.0), y + Inches(1.5), Inches(0.3), Inches(0.5),
                     "→", font_size=24, color=MID_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4),
             "ツールの操作は合計5分程度。商談の自然な流れの中でシームレスに活用できます。",
             font_size=14, color=ACCENT_AMBER, bold=True)

add_page_number(slide, 11, TOTAL_SLIDES)


# ============================================================
# SLIDE 12: ロードマップ & まとめ
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)
add_gradient_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE, ACCENT_CYAN)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(6), Inches(0.6),
             "ROADMAP & SUMMARY", font_size=13, color=ACCENT_CYAN, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.8),
             "開発ロードマップ & まとめ", font_size=36, color=WHITE, bold=True)

# ロードマップ
roadmap = [
    ("MVP（現在）", "v1.0",
     "入力UI / ROIダッシュボード / シナリオ比較\n業種テンプレート / 感度分析 / 導入比較グラフ",
     ACCENT_BLUE),
    ("次期アップデート", "v1.1",
     "AI適用領域マッピング / 導入ロードマップ\nPDF出力 / ケーススタディ / 共有リンク",
     ACCENT_CYAN),
    ("将来構想", "v2.0",
     "LLMによるエグゼクティブサマリー自動生成\nAIが試算結果を分析し考察を提供",
     ACCENT_GREEN),
]

for i, (phase, ver, features, color) in enumerate(roadmap):
    x = Inches(0.8 + i * 4.1)
    card = add_rounded_card(slide, x, Inches(1.9), Inches(3.7), Inches(2.2))
    add_gradient_bar(slide, x, Inches(1.9), Inches(3.7), Inches(0.05), color, color)

    ver_badge = add_rounded_card(slide, x + Inches(0.2), Inches(2.15), Inches(0.7), Inches(0.3), color)
    add_text_box(slide, x + Inches(0.2), Inches(2.17), Inches(0.7), Inches(0.25),
                 ver, font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x + Inches(1.1), Inches(2.15), Inches(2.4), Inches(0.35),
                 phase, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.65), Inches(3.1), Inches(1.2),
                 features, font_size=11, color=LIGHT_GRAY)

# まとめ
summary_card = add_rounded_card(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5),
                                 RGBColor(0x15, 0x20, 0x33))
add_gradient_bar(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.05), ACCENT_AMBER, ACCENT_AMBER)

add_text_box(slide, Inches(1.2), Inches(4.75), Inches(10), Inches(0.5),
             "AI-ROI Simulator が実現すること", font_size=22, color=WHITE, bold=True)

summary_items = [
    ("初回商談で即座にROI試算", "「御社の場合」という具体的な数字で経営層を動かす"),
    ("3シナリオの透明な提示", "保守的でも効果がある」と証明し、信頼を獲得する"),
    ("インタラクティブな探索", "What-if分析にリアルタイムで対応し、専門性を示す"),
    ("日本市場初のWebツール", "競合コンサルタントとの決定的な差別化を実現する"),
]

for i, (title, desc) in enumerate(summary_items):
    x = Inches(1.4) if i < 2 else Inches(7.0)
    y = Inches(5.4 + (i % 2) * 0.7)
    add_text_box(slide, x, y, Inches(0.3), Inches(0.3),
                 "✓", font_size=16, color=ACCENT_GREEN, bold=True)
    add_text_box(slide, x + Inches(0.35), y, Inches(5), Inches(0.3),
                 title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.35), y + Inches(0.3), Inches(5), Inches(0.3),
                 desc, font_size=11, color=LIGHT_GRAY)

add_page_number(slide, 12, TOTAL_SLIDES)


# ─── 保存 ───
output_path = "/home/user/AI-ROI-simulator_v1/AI-ROI-Simulator_Guide.pptx"
prs.save(output_path)
print(f"PowerPoint saved to: {output_path}")
