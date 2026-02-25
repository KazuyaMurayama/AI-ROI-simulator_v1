#!/usr/bin/env python3
"""AI-ROI-simulator_v1 v2: White-background, 20-slide presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# ─── Brand Colors (white-background palette) ───
WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x1E, 0x29, 0x3B)           # Primary text
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)      # Secondary text
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)       # Tertiary / labels
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)     # Borders / dividers
CARD_BG = RGBColor(0xF1, 0xF5, 0xF9)        # Card background
CARD_BG_BLUE = RGBColor(0xEF, 0xF6, 0xFF)   # Blue card
CARD_BG_GREEN = RGBColor(0xEC, 0xFD, 0xF5)  # Green card
CARD_BG_AMBER = RGBColor(0xFF, 0xFB, 0xEB)  # Amber card
CARD_BG_CYAN = RGBColor(0xEC, 0xFE, 0xFF)   # Cyan card

ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB)    # Deeper blue for white bg
ACCENT_CYAN = RGBColor(0x05, 0x91, 0xB3)
ACCENT_GREEN = RGBColor(0x05, 0x96, 0x69)
ACCENT_AMBER = RGBColor(0xD9, 0x77, 0x06)
ACCENT_RED = RGBColor(0xDC, 0x26, 0x26)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = Inches(13.333)

TOTAL_SLIDES = 20

# ─── Safe layout constants ───
ML = Inches(0.8)      # margin left
MR = Inches(0.8)      # margin right
CONTENT_W = Inches(11.733)  # usable width = 13.333 - 0.8*2
HEADER_Y = Inches(0.35)
TITLE_Y = Inches(0.75)
CONTENT_TOP = Inches(1.65)
FOOTER_Y = Inches(7.05)


def white_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE_BG


def accent_bar(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def card(slide, left, top, width, height, fill_color=CARD_BG, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def txt(slide, left, top, width, height, text, size=16, color=NAVY,
        bold=False, align=PP_ALIGN.LEFT, font="Meiryo"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tb


def page_num(slide, num):
    txt(slide, Inches(12.0), FOOTER_Y, Inches(1.0), Inches(0.3),
        f"{num} / {TOTAL_SLIDES}", size=9, color=MID_GRAY, align=PP_ALIGN.RIGHT)


def slide_header(slide, tag, title, tag_color=ACCENT_BLUE):
    """Standard header: top accent bar + tag + title."""
    accent_bar(slide, Inches(0), Inches(0), SW, Inches(0.05), tag_color)
    txt(slide, ML, HEADER_Y, Inches(4), Inches(0.3), tag,
        size=11, color=tag_color, bold=True)
    txt(slide, ML, TITLE_Y, Inches(11), Inches(0.7), title,
        size=32, color=NAVY, bold=True)


def bullet_list(slide, left, top, width, items, size=13, spacing=0.55):
    """Add a bulleted list. items = [(text, color_or_None), ...]."""
    for i, (text, c) in enumerate(items):
        y = top + Inches(i * spacing)
        txt(slide, left, y, Inches(0.3), Inches(0.35),
            "●", size=8, color=c if c else ACCENT_BLUE)
        txt(slide, left + Inches(0.3), y, width - Inches(0.3), Inches(0.45),
            text, size=size, color=NAVY)


# ============================================================
# 1. TITLE SLIDE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
accent_bar(sl, Inches(0), Inches(0), SW, Inches(0.08), ACCENT_BLUE)
accent_bar(sl, ML, Inches(5.6), Inches(4.5), Inches(0.04), ACCENT_BLUE)

txt(sl, ML, Inches(2.0), Inches(10), Inches(1.0),
    "AI-ROI Simulator", size=52, color=NAVY, bold=True)
txt(sl, ML, Inches(3.2), Inches(10), Inches(0.7),
    "AI導入ROI試算ツール  |  機能説明 & 使い方ガイド", size=24, color=ACCENT_BLUE)
txt(sl, ML, Inches(4.3), Inches(9), Inches(0.8),
    "初回商談でクライアントに「御社のAI導入ROI」を即座にシミュレーション。\n提案の説得力を劇的に高める、AIコンサルタント専用ツールです。",
    size=15, color=DARK_GRAY)

badge = card(sl, ML, Inches(5.8), Inches(3.5), Inches(0.55), ACCENT_BLUE)
txt(sl, ML + Inches(0.2), Inches(5.85), Inches(3.1), Inches(0.4),
    "AI Transformation Architect", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

txt(sl, Inches(9.5), Inches(6.5), Inches(3), Inches(0.3),
    "Version 1.0  |  2026", size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)
page_num(sl, 1)


# ============================================================
# 2. OVERVIEW
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "OVERVIEW", "AI-ROI Simulator とは？")

txt(sl, ML, CONTENT_TOP, CONTENT_W, Inches(0.8),
    "クライアントの業種・規模・対象業務を入力するだけで、AI導入の投資対効果（ROI）を即座に試算。\n3つのシナリオ比較とインタラクティブな感度分析で、経営層の意思決定を後押しするWebツールです。",
    size=15, color=DARK_GRAY)

items = [
    ("商談クロージングツール",
     "初回商談30分で「御社のAI導入ROI」をリアルタイムにシミュレーション。数字の力で次回アポ・受注を引き出します。",
     ACCENT_BLUE, CARD_BG_BLUE),
    ("ポートフォリオ",
     "React + TypeScript で構築されたプロフェッショナルなWebアプリ。AI変革アーキテクトとしての技術力の証明です。",
     ACCENT_CYAN, CARD_BG_CYAN),
    ("差別化の武器",
     "日本市場初のインタラクティブAI ROIシミュレーター。競合コンサルタントとの決定的な差別化要因になります。",
     ACCENT_GREEN, CARD_BG_GREEN),
]

for i, (title, desc, color, bg) in enumerate(items):
    x = ML + Inches(i * 4.0)
    w = Inches(3.7)
    y = Inches(3.0)
    card(sl, x, y, w, Inches(3.2), bg, color)
    accent_bar(sl, x, y, w, Inches(0.05), color)
    txt(sl, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.35),
        title, size=18, color=color, bold=True)
    txt(sl, x + Inches(0.3), y + Inches(0.85), Inches(3.1), Inches(2.0),
        desc, size=13, color=DARK_GRAY)

txt(sl, ML, Inches(6.4), CONTENT_W, Inches(0.4),
    "「この数字の根拠は？」に即座に答えられる — 前提条件の完全な透明性を確保したツール設計",
    size=13, color=ACCENT_AMBER, bold=True)
page_num(sl, 2)


# ============================================================
# 3. STRATEGIC POSITIONING
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "STRATEGY", "戦略的ポジショニング")

rows = [
    ("課題", "多くのAIコンサルタントが「AI導入で効率化できます」と定性的に提案。\nクライアント経営層は「で、いくら儲かるの？」という定量的な回答を求めている。",
     ACCENT_RED, Inches(1.7)),
    ("解決策", "初回商談の場で、クライアント固有の条件に基づいたROI試算を即座に提示。\n「御社の場合、年間2,400万円のコスト削減が見込めます」— この具体性が他社との違いを生む。",
     ACCENT_BLUE, Inches(3.3)),
    ("実績の裏付け", "Amazon Japan広告部門でROAS 30%改善・累計売上18億円の実績。\n5社へのAI導入支援で各社1,000万円以上の成果。これらをベンチマークに組み込み。",
     ACCENT_GREEN, Inches(4.9)),
]

for label, desc, color, y in rows:
    card(sl, ML, y, CONTENT_W, Inches(1.35), CARD_BG, color)
    accent_bar(sl, ML, y, Inches(0.06), Inches(1.35), color)
    txt(sl, ML + Inches(0.4), y + Inches(0.15), Inches(2), Inches(0.35),
        label, size=16, color=color, bold=True)
    txt(sl, ML + Inches(0.4), y + Inches(0.55), Inches(10.8), Inches(0.7),
        desc, size=13, color=DARK_GRAY)

page_num(sl, 3)


# ============================================================
# 4. FEATURE OVERVIEW (6 MVP features)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "FEATURES", "主要機能一覧（MVP 6機能）")

features = [
    ("F1", "パラメータ入力UI", "業種・規模・対象業務をステッパー形式で入力。\n商談の流れに沿ったガイド付きフォーム。", ACCENT_BLUE, CARD_BG_BLUE),
    ("F2", "ROIダッシュボード", "ヒーローナンバー + KPIカード + グラフで\n試算結果を一目で把握。30秒で理解可能。", ACCENT_CYAN, CARD_BG_CYAN),
    ("F3", "3シナリオ比較", "保守的 / 標準 / 楽観的の3シナリオを\n並列表示。信頼性と説得力の両立。", ACCENT_GREEN, CARD_BG_GREEN),
    ("F4", "業種別テンプレート", "10業種のプリセットパラメータ。\n業種選択だけで即座に試算を開始。", ACCENT_AMBER, CARD_BG_AMBER),
    ("F9", "感度分析", "スライダー操作でROIがリアルタイム変動。\nExcelでは不可能なインタラクティブ体験。", ACCENT_BLUE, CARD_BG_BLUE),
    ("F11", "導入あり/なし比較", "AI導入前後の3〜5年推移グラフ。\n経営層が最も直感的に理解できる表示。", ACCENT_CYAN, CARD_BG_CYAN),
]

for i, (fid, title, desc, color, bg) in enumerate(features):
    row, col = divmod(i, 3)
    x = ML + Inches(col * 4.0)
    y = CONTENT_TOP + Inches(row * 2.6)
    w, h = Inches(3.7), Inches(2.3)
    card(sl, x, y, w, h, bg, color)
    # badge
    badge_s = card(sl, x + Inches(0.2), y + Inches(0.2), Inches(0.55), Inches(0.3), color)
    txt(sl, x + Inches(0.2), y + Inches(0.22), Inches(0.55), Inches(0.25),
        fid, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.9), y + Inches(0.2), Inches(2.5), Inches(0.35),
        title, size=16, color=NAVY, bold=True)
    txt(sl, x + Inches(0.3), y + Inches(0.7), Inches(3.1), Inches(1.3),
        desc, size=12, color=DARK_GRAY)

page_num(sl, 4)


# ============================================================
# 5. FEATURE: F1 Parameter Input
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "FEATURE  F1", "パラメータ入力UI — 二段構え設計")

# Left: form mockup
card(sl, ML, CONTENT_TOP, Inches(5.5), Inches(4.8), CARD_BG)
txt(sl, ML + Inches(0.3), CONTENT_TOP + Inches(0.15), Inches(3), Inches(0.3),
    "入力フォーム（イメージ）", size=13, color=MID_GRAY, bold=True)

form_items = [
    ("業種を選択", "製造業 / 小売・EC / 金融 / IT / 医療 / 物流 / 不動産 / 教育 / メディア / 専門サービス"),
    ("従業員数", "〜50名 / 51-200名 / 201-500名 / 501-1,000名 / 1,001-5,000名 / 5,001名〜"),
    ("年間売上規模", "〜1億円 / 1-10億円 / 10-50億円 / 50-100億円 / 100-500億円 / 500億円〜"),
    ("AI導入対象業務", "カスタマーサポート / 営業・マーケ / データ分析 / 文書管理 / 経理 / 人事 等"),
]

for i, (label, opts) in enumerate(form_items):
    fy = CONTENT_TOP + Inches(0.6 + i * 1.05)
    txt(sl, ML + Inches(0.4), fy, Inches(4.5), Inches(0.3),
        label, size=13, color=ACCENT_BLUE, bold=True)
    card(sl, ML + Inches(0.4), fy + Inches(0.3), Inches(4.5), Inches(0.5), WHITE_BG, LIGHT_GRAY)
    txt(sl, ML + Inches(0.6), fy + Inches(0.33), Inches(4.1), Inches(0.4),
        opts, size=9, color=MID_GRAY)

# Right: explanation
rx = Inches(7.0)
txt(sl, rx, CONTENT_TOP, Inches(5.5), Inches(0.4),
    "設計のポイント", size=20, color=NAVY, bold=True)

points = [
    ("二段構え設計", "必須入力（4項目）だけで概算ROIが出ます。\n任意入力でさらに精度を上げることが可能。", ACCENT_BLUE, CARD_BG_BLUE),
    ("商談で聞ける範囲のみ", "必須項目はすべて初回商談のヒアリングで\n自然に聞ける情報に限定しています。", ACCENT_GREEN, CARD_BG_GREEN),
    ("業種テンプレート連動", "業種を選ぶと典型的なパラメータが自動で\nプリセット。入力の手間を最小化。", ACCENT_AMBER, CARD_BG_AMBER),
    ("リアルタイムプレビュー", "入力と同時にプレビューが更新。\n商談中のテンポを損なわない設計。", ACCENT_CYAN, CARD_BG_CYAN),
]

for i, (title, desc, color, bg) in enumerate(points):
    py = CONTENT_TOP + Inches(0.6 + i * 1.2)
    card(sl, rx, py, Inches(5.5), Inches(1.0), bg, color)
    accent_bar(sl, rx, py, Inches(0.05), Inches(1.0), color)
    txt(sl, rx + Inches(0.3), py + Inches(0.1), Inches(5), Inches(0.3),
        title, size=13, color=color, bold=True)
    txt(sl, rx + Inches(0.3), py + Inches(0.4), Inches(5), Inches(0.5),
        desc, size=11, color=DARK_GRAY)

page_num(sl, 5)


# ============================================================
# 6. FEATURE: F2 ROI Dashboard
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "FEATURE  F2", "ROIダッシュボード — 30秒で伝わるUI")

# Hero number mockup
card(sl, ML, CONTENT_TOP, CONTENT_W, Inches(1.5), CARD_BG_BLUE, ACCENT_BLUE)
accent_bar(sl, ML, CONTENT_TOP, CONTENT_W, Inches(0.05), ACCENT_BLUE)
txt(sl, ML + Inches(0.4), CONTENT_TOP + Inches(0.15), Inches(6), Inches(0.3),
    "御社のAI導入による年間効果（標準シナリオ）", size=13, color=DARK_GRAY)
txt(sl, ML + Inches(0.4), CONTENT_TOP + Inches(0.5), Inches(6), Inches(0.7),
    "¥ 2,400 万円", size=42, color=ACCENT_GREEN, bold=True)
txt(sl, ML + Inches(0.4), CONTENT_TOP + Inches(1.1), Inches(4), Inches(0.25),
    "年間コスト削減額", size=12, color=MID_GRAY)

# KPI cards
kpis = [
    ("3年ROI", "320%", ACCENT_BLUE, CARD_BG_BLUE),
    ("投資回収期間", "8.5 ヶ月", ACCENT_CYAN, CARD_BG_CYAN),
    ("年間売上増加", "¥1,800万", ACCENT_GREEN, CARD_BG_GREEN),
    ("FTE換算", "3.2 人分", ACCENT_AMBER, CARD_BG_AMBER),
]

for i, (label, value, color, bg) in enumerate(kpis):
    x = ML + Inches(i * 3.0)
    y = Inches(3.5)
    card(sl, x, y, Inches(2.7), Inches(1.2), bg, color)
    accent_bar(sl, x, y, Inches(2.7), Inches(0.04), color)
    txt(sl, x + Inches(0.2), y + Inches(0.15), Inches(2.3), Inches(0.25),
        label, size=11, color=MID_GRAY)
    txt(sl, x + Inches(0.2), y + Inches(0.45), Inches(2.3), Inches(0.5),
        value, size=26, color=color, bold=True)

# Explanation
txt(sl, ML, Inches(5.0), CONTENT_W, Inches(0.35),
    "「30秒ルール」— 経営層が画面を見て30秒で「うちにとっての効果」を理解できるUI設計",
    size=14, color=ACCENT_AMBER, bold=True)

desc_lines = [
    ("最上部のヒーローナンバーで最も重要な数字を大きく表示し、KPIカードで主要指標を一覧表示します。", None),
    ("すべての数字はクライアントが入力した情報に基づくパーソナライズされた試算です。", None),
    ("※ 試算値は公開調査レポートおよび類似プロジェクト実績に基づく推定値です。", None),
]
for i, (line, _) in enumerate(desc_lines):
    txt(sl, ML, Inches(5.5 + i * 0.4), CONTENT_W, Inches(0.35),
        line, size=12, color=DARK_GRAY)

page_num(sl, 6)


# ============================================================
# 7. FEATURE: F3 Scenario Comparison
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "FEATURE  F3", "3シナリオ比較 — 信頼性の防御線")

scenarios = [
    ("保守的", "Conservative",
     [("年間効果", "¥1,200万"), ("3年ROI", "180%"), ("回収期間", "14ヶ月")],
     MID_GRAY, CARD_BG, "最低限期待できる効果。\nすべてのパラメータに\n下限値を適用。"),
    ("標準", "Standard",
     [("年間効果", "¥2,400万"), ("3年ROI", "320%"), ("回収期間", "8.5ヶ月")],
     ACCENT_BLUE, CARD_BG_BLUE, "一般的に期待できる効果。\n業界平均の\nベンチマーク値を適用。"),
    ("楽観的", "Optimistic",
     [("年間効果", "¥4,100万"), ("3年ROI", "580%"), ("回収期間", "5ヶ月")],
     ACCENT_GREEN, CARD_BG_GREEN, "最大限の効果が出た場合。\nすべてのパラメータに\n上限値を適用。"),
]

for i, (name, eng, kpis_list, color, bg, desc) in enumerate(scenarios):
    x = ML + Inches(i * 4.0)
    w = Inches(3.7)
    card(sl, x, CONTENT_TOP, w, Inches(4.5), bg, color)
    accent_bar(sl, x, CONTENT_TOP, w, Inches(0.05), color)
    txt(sl, x + Inches(0.3), CONTENT_TOP + Inches(0.2), Inches(3.1), Inches(0.4),
        name, size=22, color=color, bold=True)
    txt(sl, x + Inches(0.3), CONTENT_TOP + Inches(0.6), Inches(3.1), Inches(0.25),
        eng, size=10, color=MID_GRAY)

    for j, (kl, kv) in enumerate(kpis_list):
        ky = CONTENT_TOP + Inches(1.1 + j * 0.65)
        txt(sl, x + Inches(0.3), ky, Inches(1.5), Inches(0.25), kl, size=11, color=MID_GRAY)
        txt(sl, x + Inches(2.0), ky, Inches(1.4), Inches(0.25),
            kv, size=15, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)

    txt(sl, x + Inches(0.3), CONTENT_TOP + Inches(3.2), Inches(3.1), Inches(1.0),
        desc, size=11, color=DARK_GRAY)

txt(sl, ML, Inches(6.4), CONTENT_W, Inches(0.4),
    "「保守的に見積もっても年間1,200万円の効果が見込めます」— この一言がクライアントの信頼を勝ち取る",
    size=14, color=ACCENT_AMBER, bold=True)
page_num(sl, 7)


# ============================================================
# 8. FEATURE: F9 Sensitivity Analysis
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "FEATURE  F9", "感度分析 — Excelとの決定的な違い")

# Left: slider mockup
card(sl, ML, CONTENT_TOP, Inches(5.5), Inches(4.5), CARD_BG)
txt(sl, ML + Inches(0.3), CONTENT_TOP + Inches(0.15), Inches(3), Inches(0.3),
    "スライダー操作（イメージ）", size=13, color=MID_GRAY, bold=True)

sliders = [
    ("AI自動化率", "25%", 0.5, ACCENT_BLUE),
    ("導入投資額", "¥2,000万", 0.4, ACCENT_CYAN),
    ("効果発現期間", "6ヶ月", 0.35, ACCENT_GREEN),
    ("年間運用費率", "25%", 0.5, ACCENT_AMBER),
]

for i, (label, value, ratio, color) in enumerate(sliders):
    sy = CONTENT_TOP + Inches(0.6 + i * 1.0)
    txt(sl, ML + Inches(0.4), sy, Inches(2.5), Inches(0.25),
        label, size=13, color=NAVY, bold=True)
    txt(sl, ML + Inches(3.5), sy, Inches(1.8), Inches(0.25),
        value, size=13, color=color, bold=True, align=PP_ALIGN.RIGHT)
    # slider track
    card(sl, ML + Inches(0.4), sy + Inches(0.32), Inches(4.5), Inches(0.12), LIGHT_GRAY)
    # slider fill
    card(sl, ML + Inches(0.4), sy + Inches(0.32), Inches(4.5 * ratio), Inches(0.12), color)

# Right: benefits
rx = Inches(7.0)
txt(sl, rx, CONTENT_TOP, Inches(5.5), Inches(0.4),
    "この機能の価値", size=20, color=NAVY, bold=True)

benefits = [
    ("リアルタイム更新", "スライダーを動かすと、グラフとKPIが\n即座に更新されます。", ACCENT_BLUE),
    ("What-if対応", "「もし自動化率がもう少し低かったら？」\nその場で答えを見せられます。", ACCENT_GREEN),
    ("影響度の可視化", "各パラメータがROIに与える影響度を\nトルネードチャートで表示。", ACCENT_CYAN),
    ("Excel操作不要", "クリック＆ドラッグだけで探索完了。\n商談中にExcelを開く必要なし。", ACCENT_AMBER),
]

for i, (title, desc, color) in enumerate(benefits):
    by = CONTENT_TOP + Inches(0.6 + i * 1.15)
    card(sl, rx, by, Inches(5.5), Inches(0.95), CARD_BG, color)
    accent_bar(sl, rx, by, Inches(0.05), Inches(0.95), color)
    txt(sl, rx + Inches(0.3), by + Inches(0.1), Inches(5), Inches(0.3),
        title, size=13, color=color, bold=True)
    txt(sl, rx + Inches(0.3), by + Inches(0.4), Inches(5), Inches(0.45),
        desc, size=11, color=DARK_GRAY)

page_num(sl, 8)


# ============================================================
# 9. FEATURE: F11 Before/After
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "FEATURE  F11", "導入あり/なし比較 — 5年推移グラフ")

chart_data = CategoryChartData()
chart_data.categories = ['1年目', '2年目', '3年目', '4年目', '5年目']
chart_data.add_series('AI導入なし（コスト推移）', (100, 102, 104, 106, 108))
chart_data.add_series('AI導入あり（コスト推移）', (85, 78, 72, 68, 65))
chart_data.add_series('累計コスト削減効果', (15, 39, 71, 109, 152))

cf = sl.shapes.add_chart(
    XL_CHART_TYPE.LINE_MARKERS, ML, CONTENT_TOP, Inches(7.5), Inches(4.2), chart_data
)
chart = cf.chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = DARK_GRAY

plot = chart.plots[0]
series_colors = [MID_GRAY, ACCENT_BLUE, ACCENT_GREEN]
for i, series in enumerate(plot.series):
    series.format.line.color.rgb = series_colors[i]
    series.format.line.width = Pt(3)
    series.smooth = True

# Right side
rx = Inches(8.8)
txt(sl, rx, CONTENT_TOP, Inches(4), Inches(0.35),
    "グラフの読み方", size=18, color=NAVY, bold=True)

legends = [
    ("グレーの線", "AI未導入の場合のコスト推移\n（現状維持シナリオ）", MID_GRAY),
    ("ブルーの線", "AI導入後のコスト推移\n（削減効果を反映）", ACCENT_BLUE),
    ("グリーンの線", "累計コスト削減効果\n（差額の積み上げ）", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(legends):
    ly = CONTENT_TOP + Inches(0.6 + i * 1.1)
    accent_bar(sl, rx, ly, Inches(0.5), Inches(0.04), color)
    txt(sl, rx, ly + Inches(0.15), Inches(3.8), Inches(0.25),
        title, size=12, color=color, bold=True)
    txt(sl, rx, ly + Inches(0.4), Inches(3.8), Inches(0.5),
        desc, size=11, color=DARK_GRAY)

txt(sl, rx, Inches(5.0), Inches(4), Inches(0.7),
    "経営層が最も直感的に理解できる\n「Before / After」の可視化です。",
    size=13, color=ACCENT_AMBER, bold=True)
page_num(sl, 9)


# ============================================================
# 10. FEATURE: F4 Industry Templates
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "FEATURE  F4", "業種別テンプレート — ワンクリックで試算開始")

txt(sl, ML, CONTENT_TOP, CONTENT_W, Inches(0.5),
    "10業種に対応したプリセットパラメータを搭載。業種を選ぶだけで、その業種に典型的な\nAI導入シナリオが自動的にセットされ、すぐに試算を開始できます。",
    size=14, color=DARK_GRAY)

industries = [
    ("製造業", "品質管理、SCM最適化、\n予知保全", ACCENT_BLUE),
    ("小売・EC", "需要予測、パーソナライズ、\n在庫最適化", ACCENT_GREEN),
    ("金融・保険", "審査自動化、不正検知、\n事務処理効率化", ACCENT_CYAN),
    ("IT・通信", "コード生成支援、\nテスト自動化", ACCENT_AMBER),
    ("医療・HC", "画像診断支援、\n文書作成自動化", ACCENT_BLUE),
    ("物流・運輸", "ルート最適化、\n配車計画、倉庫管理", ACCENT_GREEN),
    ("不動産", "物件査定自動化、\n営業効率化", ACCENT_CYAN),
    ("教育", "コンテンツ生成、\n個別指導最適化", ACCENT_AMBER),
    ("メディア・広告", "広告最適化、\nクリエイティブ生成", ACCENT_BLUE),
    ("専門サービス", "ナレッジ管理、\nリサーチ自動化", ACCENT_GREEN),
]

for i, (name, desc, color) in enumerate(industries):
    row, col = divmod(i, 5)
    x = ML + Inches(col * 2.4)
    y = Inches(2.7 + row * 2.0)
    card(sl, x, y, Inches(2.2), Inches(1.7), CARD_BG, color)
    accent_bar(sl, x, y, Inches(2.2), Inches(0.04), color)
    txt(sl, x + Inches(0.2), y + Inches(0.15), Inches(1.8), Inches(0.3),
        name, size=14, color=color, bold=True)
    txt(sl, x + Inches(0.2), y + Inches(0.55), Inches(1.8), Inches(1.0),
        desc, size=10, color=DARK_GRAY)

page_num(sl, 10)


# ============================================================
# 11. HOW TO USE: 5-Step Overview
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "HOW TO USE", "使い方：5ステップで即座にROI試算")

steps = [
    ("Step 1", "業種・規模を入力", "業種、従業員数、売上規模、\n対象業務を選択", ACCENT_BLUE),
    ("Step 2", "ROI結果を確認", "ヒーローナンバーと\nKPIカードで即座に把握", ACCENT_CYAN),
    ("Step 3", "シナリオ比較", "保守的/標準/楽観的の\n3パターンを比較", ACCENT_GREEN),
    ("Step 4", "感度分析", "スライダーでパラメータを\n調整し探索", ACCENT_AMBER),
    ("Step 5", "将来推移を確認", "3〜5年の導入効果推移を\nグラフで提示", ACCENT_BLUE),
]

for i, (step, title, desc, color) in enumerate(steps):
    x = ML + Inches(i * 2.45)
    y = CONTENT_TOP
    w = Inches(2.2)

    # Step number badge
    card(sl, x + Inches(0.5), y, Inches(1.2), Inches(0.4), color)
    txt(sl, x + Inches(0.5), y + Inches(0.05), Inches(1.2), Inches(0.3),
        step, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Title
    txt(sl, x, y + Inches(0.6), w, Inches(0.35),
        title, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Card
    card(sl, x, y + Inches(1.1), w, Inches(1.5), CARD_BG, color)
    txt(sl, x + Inches(0.2), y + Inches(1.25), Inches(1.8), Inches(1.2),
        desc, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)

    # Arrow
    if i < 4:
        txt(sl, x + Inches(2.2), y + Inches(1.5), Inches(0.3), Inches(0.3),
            "→", size=20, color=MID_GRAY)

# Bottom flow summary
card(sl, ML, Inches(4.5), CONTENT_W, Inches(2.2), CARD_BG_BLUE, ACCENT_BLUE)
txt(sl, ML + Inches(0.4), Inches(4.7), Inches(10), Inches(0.35),
    "商談での操作は合計5分程度", size=16, color=ACCENT_BLUE, bold=True)
txt(sl, ML + Inches(0.4), Inches(5.15), Inches(10), Inches(1.2),
    "1. クライアントへのヒアリングと並行して業種・規模を選択（1分）\n"
    "2. 対象業務を選択し、ROIダッシュボードの結果を見せる（1分）\n"
    "3. 3シナリオを比較し、保守的シナリオの数字を強調（1分）\n"
    "4. 「もし○○だったら？」の質問にスライダーで即答（1分）\n"
    "5. 5年推移グラフで長期的な効果を提示し、クロージングへ（1分）",
    size=12, color=DARK_GRAY)

page_num(sl, 11)


# ============================================================
# 12. USE CASE 1: Initial Meeting
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "USE CASE 1", "活用シーン：初回オンライン商談（30分）", ACCENT_BLUE)

# Timeline
phases = [
    ("0-5分", "ヒアリング",
     "自己紹介と課題のヒアリング。\nツールはまだ使わず、\nクライアントの話に集中。",
     "「御社ではAIの導入を\nご検討されていますか？」",
     MID_GRAY),
    ("5-15分", "課題特定",
     "業種を選択し、AI適用可能な\n業務領域を一緒に特定。\nテンプレートで即座にセット。",
     "「御社の業種では、この業務に\nAIが最も効果的です」",
     ACCENT_BLUE),
    ("15-25分", "ROI提示",
     "パラメータを入力し、\nリアルタイムでROI試算。\n感度分析で深堀り。",
     "「御社の場合、年間2,400万円の\n効果が見込めます」",
     ACCENT_GREEN),
    ("25-30分", "クロージング",
     "3シナリオ比較と将来推移を\n提示。次回アポイントの取得。",
     "「詳細なレポートを\nお送りいたします」",
     ACCENT_AMBER),
]

for i, (time, title, action, script, color) in enumerate(phases):
    x = ML + Inches(i * 3.1)
    y = CONTENT_TOP
    w = Inches(2.8)

    # Time badge
    card(sl, x + Inches(0.4), y, Inches(1.4), Inches(0.35), color)
    txt(sl, x + Inches(0.4), y + Inches(0.03), Inches(1.4), Inches(0.3),
        time, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Title
    txt(sl, x, y + Inches(0.5), w, Inches(0.3),
        title, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)

    # Action card
    card(sl, x, y + Inches(0.9), w, Inches(1.6), CARD_BG)
    txt(sl, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(1.3),
        action, size=11, color=DARK_GRAY)

    # Script
    card(sl, x, y + Inches(2.7), w, Inches(1.1), CARD_BG_BLUE)
    txt(sl, x + Inches(0.15), y + Inches(2.75), Inches(2.5), Inches(0.2),
        "トークスクリプト:", size=9, color=MID_GRAY, bold=True)
    txt(sl, x + Inches(0.15), y + Inches(2.95), Inches(2.5), Inches(0.75),
        script, size=10, color=ACCENT_BLUE)

    # Arrow
    if i < 3:
        txt(sl, x + w, y + Inches(1.5), Inches(0.3), Inches(0.3),
            "→", size=18, color=MID_GRAY)

txt(sl, ML, Inches(6.4), CONTENT_W, Inches(0.35),
    "ツール操作は合計5分。商談の自然な流れの中でシームレスに活用でき、プレゼン感を出さない。",
    size=13, color=ACCENT_AMBER, bold=True)
page_num(sl, 12)


# ============================================================
# 13. USE CASE 2: Proposal Embedding
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "USE CASE 2", "活用シーン：提案書（PPT）への埋め込み", ACCENT_CYAN)

txt(sl, ML, CONTENT_TOP, CONTENT_W, Inches(0.6),
    "商談後にクライアントに送る提案書に、ROI試算結果のスクリーンショットを埋め込みます。\n社内稟議にそのまま使える品質のビジュアルを提供します。",
    size=14, color=DARK_GRAY)

# Left: workflow
card(sl, ML, Inches(2.6), Inches(5.5), Inches(4.0), CARD_BG)
txt(sl, ML + Inches(0.3), Inches(2.75), Inches(4), Inches(0.3),
    "提案書作成ワークフロー", size=16, color=NAVY, bold=True)

wf_steps = [
    ("1.", "商談中にツールでROI試算を実施"),
    ("2.", "ダッシュボード画面をスクリーンショット"),
    ("3.", "3シナリオ比較のスクリーンショットも取得"),
    ("4.", "提案書のROIセクションにスクリーンショットを貼付"),
    ("5.", "前提条件と免責事項も併記して信頼性を担保"),
]

for i, (num, step) in enumerate(wf_steps):
    wy = Inches(3.3 + i * 0.55)
    txt(sl, ML + Inches(0.5), wy, Inches(0.4), Inches(0.3),
        num, size=14, color=ACCENT_BLUE, bold=True)
    txt(sl, ML + Inches(0.9), wy, Inches(4), Inches(0.35),
        step, size=13, color=DARK_GRAY)

# Right: benefits
rx = Inches(7.0)
txt(sl, rx, Inches(2.6), Inches(5.5), Inches(0.3),
    "この活用シーンのメリット", size=16, color=NAVY, bold=True)

benefits = [
    ("提案書作成時間 50%短縮", "数字とグラフが自動生成されるため、\nExcelでの試算作業が不要に。", ACCENT_GREEN),
    ("社内稟議に耐える品質", "プロフェッショナルなデザインで、\nクライアント社内の意思決定者にも訴求。", ACCENT_BLUE),
    ("根拠の明確化", "「この数字の根拠は？」に\n前提条件とデータソースで即答。", ACCENT_AMBER),
    ("再現性", "同じパラメータで何度でも\n同じ結果を再現。クライアントも確認可能。", ACCENT_CYAN),
]

for i, (title, desc, color) in enumerate(benefits):
    by = Inches(3.1 + i * 1.0)
    card(sl, rx, by, Inches(5.5), Inches(0.85), CARD_BG, color)
    accent_bar(sl, rx, by, Inches(0.05), Inches(0.85), color)
    txt(sl, rx + Inches(0.3), by + Inches(0.1), Inches(5), Inches(0.25),
        title, size=12, color=color, bold=True)
    txt(sl, rx + Inches(0.3), by + Inches(0.35), Inches(5), Inches(0.4),
        desc, size=11, color=DARK_GRAY)

page_num(sl, 13)


# ============================================================
# 14. USE CASE 3: Client Self-Service
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "USE CASE 3", "活用シーン：クライアントのセルフ利用", ACCENT_GREEN)

txt(sl, ML, CONTENT_TOP, CONTENT_W, Inches(0.6),
    "商談後、クライアント自身がURLにアクセスし、自社のパラメータを入力して試算を体験。\n「自分ごと」として検討を深めてもらうフェーズです。",
    size=14, color=DARK_GRAY)

# 3-column flow
flows = [
    ("商談後のフォロー", "商談後にURLを共有。\nクライアントが自分の\nペースで試算を探索。\n\n「ぜひご自身でも\nパラメータを変えて\n試してみてください」",
     ACCENT_BLUE, CARD_BG_BLUE),
    ("社内検討の材料", "クライアントが社内の\n関係者に画面を共有。\n経営会議の資料として\n活用される可能性。\n\n→ 社内の味方を増やす",
     ACCENT_GREEN, CARD_BG_GREEN),
    ("次回商談への布石", "クライアントが自ら\n試算した結果をもとに\n「もっと詳しく聞きたい」\nと次回アポにつながる。\n\n→ 能動的な受注につなげる",
     ACCENT_AMBER, CARD_BG_AMBER),
]

for i, (title, desc, color, bg) in enumerate(flows):
    x = ML + Inches(i * 4.0)
    w = Inches(3.7)
    card(sl, x, Inches(2.6), w, Inches(4.0), bg, color)
    accent_bar(sl, x, Inches(2.6), w, Inches(0.04), color)
    txt(sl, x + Inches(0.3), Inches(2.8), Inches(3.1), Inches(0.35),
        title, size=16, color=color, bold=True)
    txt(sl, x + Inches(0.3), Inches(3.3), Inches(3.1), Inches(3.0),
        desc, size=12, color=DARK_GRAY)

page_num(sl, 14)


# ============================================================
# 15. USE CASE 4: Post-Contract Baseline
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "USE CASE 4", "活用シーン：受注後の効果測定ベースライン", ACCENT_AMBER)

txt(sl, ML, CONTENT_TOP, CONTENT_W, Inches(0.6),
    "受注後のプロジェクトにおいて、ROI試算の結果をベースラインとして活用。\n「最初に試算した数字」と「実際の効果」を比較することで、プロジェクトの価値を可視化します。",
    size=14, color=DARK_GRAY)

# 2-column layout
# Left
card(sl, ML, Inches(2.6), Inches(5.7), Inches(3.8), CARD_BG, ACCENT_BLUE)
txt(sl, ML + Inches(0.3), Inches(2.8), Inches(5), Inches(0.3),
    "ベースラインとしての活用", size=16, color=ACCENT_BLUE, bold=True)

baseline_items = [
    "プロジェクト開始時に「標準シナリオ」の試算値をKPIとして設定",
    "月次・四半期ごとに実績値と試算値を比較",
    "試算値を上回った場合 → 成果をアピールし継続案件へ",
    "試算値を下回った場合 → 要因分析と改善提案で信頼を維持",
    "最終報告書に「試算 vs 実績」の比較グラフを掲載",
]
for i, item in enumerate(baseline_items):
    iy = Inches(3.3 + i * 0.55)
    txt(sl, ML + Inches(0.5), iy, Inches(0.3), Inches(0.3),
        "●", size=8, color=ACCENT_BLUE)
    txt(sl, ML + Inches(0.8), iy, Inches(4.5), Inches(0.4),
        item, size=12, color=DARK_GRAY)

# Right
rx = Inches(7.0)
card(sl, rx, Inches(2.6), Inches(5.5), Inches(3.8), CARD_BG_AMBER, ACCENT_AMBER)
txt(sl, rx + Inches(0.3), Inches(2.8), Inches(5), Inches(0.3),
    "継続案件につなげる効果", size=16, color=ACCENT_AMBER, bold=True)

effects = [
    ("信頼性の証明", "「言った通りの効果が出た」=\n次の案件の最強の営業材料"),
    ("追加提案の根拠", "効果が出た業務を起点に、\n他業務へのAI展開を提案"),
    ("長期リテーナー化", "効果測定を継続するために\n月額顧問契約につなげる"),
]
for i, (title, desc) in enumerate(effects):
    ey = Inches(3.3 + i * 1.1)
    txt(sl, rx + Inches(0.3), ey, Inches(5), Inches(0.25),
        title, size=14, color=ACCENT_AMBER, bold=True)
    txt(sl, rx + Inches(0.3), ey + Inches(0.3), Inches(5), Inches(0.6),
        desc, size=12, color=DARK_GRAY)

page_num(sl, 15)


# ============================================================
# 16. ROI MODEL
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "ROI MODEL", "ROI計算モデルの仕組み")

models = [
    ("Model A", "コスト削減型",
     "対象業務コスト × AI自動化率\n× 業種補正係数",
     "バックオフィス業務の自動化、\n問い合わせ対応の自動化",
     ACCENT_BLUE, CARD_BG_BLUE),
    ("Model B", "売上増加型",
     "現在売上 × 対象業務売上貢献率\n× AI効果係数",
     "営業最適化、広告最適化、\n需要予測",
     ACCENT_GREEN, CARD_BG_GREEN),
    ("Model C", "生産性向上型",
     "削減時間 × FTE換算\n× 平均人件費",
     "すべての業務に適用可能な\n汎用モデル（FTE換算）",
     ACCENT_CYAN, CARD_BG_CYAN),
]

for i, (mid, name, formula, use_case, color, bg) in enumerate(models):
    x = ML + Inches(i * 4.0)
    w = Inches(3.7)
    card(sl, x, CONTENT_TOP, w, Inches(3.0), bg, color)
    accent_bar(sl, x, CONTENT_TOP, w, Inches(0.04), color)
    txt(sl, x + Inches(0.3), CONTENT_TOP + Inches(0.2), Inches(3.1), Inches(0.3),
        f"{mid}: {name}", size=16, color=color, bold=True)
    # formula box
    card(sl, x + Inches(0.2), CONTENT_TOP + Inches(0.65), Inches(3.3), Inches(0.8), CARD_BG)
    txt(sl, x + Inches(0.35), CONTENT_TOP + Inches(0.7), Inches(3.0), Inches(0.7),
        formula, size=11, color=ACCENT_BLUE, font="Consolas")
    txt(sl, x + Inches(0.3), CONTENT_TOP + Inches(1.65), Inches(3.1), Inches(0.2),
        "適用場面:", size=10, color=MID_GRAY, bold=True)
    txt(sl, x + Inches(0.3), CONTENT_TOP + Inches(1.9), Inches(3.1), Inches(0.8),
        use_case, size=11, color=DARK_GRAY)

# Bottom: integrated formula
card(sl, ML, Inches(5.0), CONTENT_W, Inches(1.5), CARD_BG)
accent_bar(sl, ML, Inches(5.0), CONTENT_W, Inches(0.04), ACCENT_AMBER)
txt(sl, ML + Inches(0.4), Inches(5.2), Inches(4), Inches(0.3),
    "複合ROI計算", size=16, color=NAVY, bold=True)
txt(sl, ML + Inches(0.4), Inches(5.55), Inches(10), Inches(0.3),
    "総合年間効果 = Σ（各対象業務の [コスト削減額 + 売上増加額 + 生産性向上額]）",
    size=11, color=ACCENT_BLUE, font="Consolas")
txt(sl, ML + Inches(0.4), Inches(5.85), Inches(10), Inches(0.3),
    "3年ROI = (総合年間効果 × 3 - AI導入投資合計) / AI導入投資合計 × 100",
    size=11, color=ACCENT_BLUE, font="Consolas")
txt(sl, ML + Inches(0.4), Inches(6.15), Inches(10), Inches(0.25),
    "対象業務ごとに最適なモデルを自動選択し、重複を排除した上で総合効果を算出します。",
    size=11, color=MID_GRAY)

page_num(sl, 16)


# ============================================================
# 17. DATA TRUST & SOURCES
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "TRUST & TRANSPARENCY", "データの信頼性と根拠")

# Table
sources = [
    ("McKinsey Global AI Survey", "業種別AI導入効果の平均値", "★★★★★"),
    ("Accenture AI Profit Impact", "業種別利益率改善の予測", "★★★★☆"),
    ("MIT Sloan Management Review", "企業規模別AI導入ROI", "★★★★☆"),
    ("IDC / Gartner Reports", "IT投資規模のベンチマーク", "★★★★★"),
    ("Deloitte AI Institute", "業種別AI活用事例・効果", "★★★★☆"),
    ("自社5社AI導入実績", "特定領域のキャリブレーション", "★★★☆☆"),
]

# Header
card(sl, ML, CONTENT_TOP, Inches(7.5), Inches(0.4), ACCENT_BLUE)
txt(sl, ML + Inches(0.3), CONTENT_TOP + Inches(0.05), Inches(3), Inches(0.3),
    "データソース", size=11, color=WHITE, bold=True)
txt(sl, ML + Inches(3.3), CONTENT_TOP + Inches(0.05), Inches(2.5), Inches(0.3),
    "用途", size=11, color=WHITE, bold=True)
txt(sl, ML + Inches(6.2), CONTENT_TOP + Inches(0.05), Inches(1.2), Inches(0.3),
    "信頼度", size=11, color=WHITE, bold=True)

for i, (source, usage, rating) in enumerate(sources):
    ry = CONTENT_TOP + Inches(0.45 + i * 0.42)
    bg = CARD_BG if i % 2 == 0 else WHITE_BG
    card(sl, ML, ry, Inches(7.5), Inches(0.38), bg)
    txt(sl, ML + Inches(0.3), ry + Inches(0.05), Inches(3), Inches(0.25),
        source, size=10, color=NAVY)
    txt(sl, ML + Inches(3.3), ry + Inches(0.05), Inches(2.5), Inches(0.25),
        usage, size=10, color=DARK_GRAY)
    txt(sl, ML + Inches(6.2), ry + Inches(0.05), Inches(1.2), Inches(0.25),
        rating, size=10, color=ACCENT_AMBER)

# Right: trust design
rx = Inches(8.8)
txt(sl, rx, CONTENT_TOP, Inches(4), Inches(0.3),
    "信頼性を担保する4つの設計", size=16, color=NAVY, bold=True)

trust_pts = [
    ("3シナリオ比較", "楽観的な数字だけでなく、保守的\nシナリオを必ず並列で提示。", ACCENT_BLUE),
    ("前提条件の明示", "すべての試算結果に計算の\n前提条件を表示。", ACCENT_GREEN),
    ("データソース開示", "ベンチマーク値の出典を\nいつでも確認可能。", ACCENT_CYAN),
    ("免責事項の表示", "推定値であることを明記。\nコンサルとしての誠実さ。", ACCENT_AMBER),
]

for i, (title, desc, color) in enumerate(trust_pts):
    ty = CONTENT_TOP + Inches(0.5 + i * 1.15)
    card(sl, rx, ty, Inches(4), Inches(0.95), CARD_BG, color)
    accent_bar(sl, rx, ty, Inches(0.04), Inches(0.95), color)
    txt(sl, rx + Inches(0.3), ty + Inches(0.1), Inches(3.5), Inches(0.25),
        title, size=12, color=color, bold=True)
    txt(sl, rx + Inches(0.3), ty + Inches(0.38), Inches(3.5), Inches(0.45),
        desc, size=10, color=DARK_GRAY)

page_num(sl, 17)


# ============================================================
# 18. COMPETITIVE ADVANTAGES
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "DIFFERENTIATION", "競合との差別化ポイント")

txt(sl, ML, CONTENT_TOP, CONTENT_W, Inches(0.5),
    "市場調査の結果、日本語のインタラクティブAI ROIシミュレーターは存在しません。\n以下の8つの観点で、既存ツール・競合コンサルタントとの明確な差別化を実現します。",
    size=14, color=DARK_GRAY)

diffs = [
    ("マルチ業種×マルチ業務", "10業種×10業務カテゴリの組み合わせ対応", ACCENT_BLUE),
    ("3シナリオ比較", "既存ツールのほとんどが単一シナリオのみ", ACCENT_GREEN),
    ("インタラクティブ感度分析", "スライダーでリアルタイム変動する体験", ACCENT_CYAN),
    ("実績ベースのキャリブレーション", "5社のAI導入実績による精度補正", ACCENT_AMBER),
    ("コンサル営業に最適化されたUX", "商談フローに沿ったステッパーUI設計", ACCENT_BLUE),
    ("プレゼンテーション品質ビジュアル", "画面共有・スクショで映えるデザイン", ACCENT_GREEN),
    ("日本市場初のインタラクティブツール", "日本語対応のWebツールは競合不在", ACCENT_CYAN),
    ("前提条件の完全な透明性", "根拠とデータソースを全面開示", ACCENT_AMBER),
]

for i, (title, desc, color) in enumerate(diffs):
    row, col = divmod(i, 2)
    x = ML + Inches(col * 6.1)
    y = Inches(2.5 + row * 1.05)
    w = Inches(5.7)

    card(sl, x, y, w, Inches(0.85), CARD_BG, color)
    accent_bar(sl, x, y, Inches(0.05), Inches(0.85), color)

    num_badge = card(sl, x + Inches(0.2), y + Inches(0.18), Inches(0.4), Inches(0.4), color)
    txt(sl, x + Inches(0.2), y + Inches(0.22), Inches(0.4), Inches(0.3),
        str(i+1), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.75), y + Inches(0.12), Inches(4.7), Inches(0.3),
        title, size=13, color=NAVY, bold=True)
    txt(sl, x + Inches(0.75), y + Inches(0.42), Inches(4.7), Inches(0.3),
        desc, size=11, color=DARK_GRAY)

page_num(sl, 18)


# ============================================================
# 19. PRACTICAL TIPS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "TIPS", "商談を成功させるための実践テクニック")

# Left column
txt(sl, ML, CONTENT_TOP, Inches(5.5), Inches(0.35),
    "商談前の準備", size=18, color=NAVY, bold=True)

prep_tips = [
    ("事前リサーチ", "クライアントの業種・従業員数・売上規模を\n事前に調べ、テンプレートをプリセットしておく。", ACCENT_BLUE),
    ("シナリオの想定", "クライアントが関心を持ちそうな業務領域を\n2-3個想定し、入力の流れをシミュレーション。", ACCENT_GREEN),
    ("ネットワーク確認", "商談前にツールの動作確認。\n万一に備えスクリーンショットも準備。", ACCENT_AMBER),
]

for i, (title, desc, color) in enumerate(prep_tips):
    py = CONTENT_TOP + Inches(0.5 + i * 1.3)
    card(sl, ML, py, Inches(5.5), Inches(1.1), CARD_BG, color)
    accent_bar(sl, ML, py, Inches(0.05), Inches(1.1), color)
    txt(sl, ML + Inches(0.3), py + Inches(0.1), Inches(5), Inches(0.25),
        title, size=13, color=color, bold=True)
    txt(sl, ML + Inches(0.3), py + Inches(0.4), Inches(5), Inches(0.55),
        desc, size=11, color=DARK_GRAY)

# Right column
rx = Inches(7.0)
txt(sl, rx, CONTENT_TOP, Inches(5.5), Inches(0.35),
    "商談中のコツ", size=18, color=NAVY, bold=True)

meeting_tips = [
    ("保守的シナリオから見せる", "最初に保守的シナリオを提示し、\n「最低限でもこれだけ」と信頼を構築してから\n標準シナリオを見せる。", ACCENT_BLUE),
    ("クライアントに操作させる", "可能であればスライダーをクライアントに\n操作してもらう。「自分ごと」感が格段に上がる。", ACCENT_GREEN),
    ("前提条件を積極的に説明", "「この数字の前提は〜です」と自分から\n説明することで、誠実さと専門性をアピール。", ACCENT_AMBER),
]

for i, (title, desc, color) in enumerate(meeting_tips):
    my = CONTENT_TOP + Inches(0.5 + i * 1.3)
    card(sl, rx, my, Inches(5.5), Inches(1.1), CARD_BG, color)
    accent_bar(sl, rx, my, Inches(0.05), Inches(1.1), color)
    txt(sl, rx + Inches(0.3), my + Inches(0.1), Inches(5), Inches(0.25),
        title, size=13, color=color, bold=True)
    txt(sl, rx + Inches(0.3), my + Inches(0.4), Inches(5), Inches(0.55),
        desc, size=11, color=DARK_GRAY)

page_num(sl, 19)


# ============================================================
# 20. ROADMAP & SUMMARY
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
white_bg(sl)
slide_header(sl, "ROADMAP & SUMMARY", "開発ロードマップ & まとめ")

roadmap = [
    ("MVP（現在）", "v1.0",
     "入力UI / ROIダッシュボード\nシナリオ比較 / 業種テンプレート\n感度分析 / 導入比較グラフ",
     ACCENT_BLUE, CARD_BG_BLUE),
    ("次期アップデート", "v1.1",
     "AI適用領域マッピング\nPDF出力 / ケーススタディ\n共有リンク / ロードマップ生成",
     ACCENT_CYAN, CARD_BG_CYAN),
    ("将来構想", "v2.0",
     "LLMによるサマリー自動生成\nAIが試算結果を分析し考察提供",
     ACCENT_GREEN, CARD_BG_GREEN),
]

for i, (phase, ver, features, color, bg) in enumerate(roadmap):
    x = ML + Inches(i * 4.0)
    w = Inches(3.7)
    card(sl, x, CONTENT_TOP, w, Inches(2.0), bg, color)
    accent_bar(sl, x, CONTENT_TOP, w, Inches(0.04), color)
    ver_badge = card(sl, x + Inches(0.2), CONTENT_TOP + Inches(0.15), Inches(0.6), Inches(0.25), color)
    txt(sl, x + Inches(0.2), CONTENT_TOP + Inches(0.17), Inches(0.6), Inches(0.2),
        ver, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(1.0), CONTENT_TOP + Inches(0.15), Inches(2.5), Inches(0.3),
        phase, size=14, color=NAVY, bold=True)
    txt(sl, x + Inches(0.3), CONTENT_TOP + Inches(0.55), Inches(3.1), Inches(1.2),
        features, size=11, color=DARK_GRAY)

# Summary box
card(sl, ML, Inches(4.1), CONTENT_W, Inches(2.8), CARD_BG, ACCENT_BLUE)
accent_bar(sl, ML, Inches(4.1), CONTENT_W, Inches(0.04), ACCENT_BLUE)
txt(sl, ML + Inches(0.4), Inches(4.3), Inches(10), Inches(0.35),
    "AI-ROI Simulator が実現すること", size=18, color=NAVY, bold=True)

summary_items = [
    ("初回商談で即座にROI試算", "「御社の場合」という具体的な数字で経営層を動かす", ACCENT_BLUE),
    ("3シナリオの透明な提示", "「保守的でも効果がある」と証明し、信頼を獲得する", ACCENT_GREEN),
    ("インタラクティブな探索", "What-if分析にリアルタイムで対応し、専門性を示す", ACCENT_CYAN),
    ("日本市場初のWebツール", "競合コンサルタントとの決定的な差別化を実現する", ACCENT_AMBER),
]

for i, (title, desc, color) in enumerate(summary_items):
    col = 0 if i < 2 else 1
    row = i % 2
    sx = ML + Inches(0.4 + col * 5.8)
    sy = Inches(4.9 + row * 0.9)
    txt(sl, sx, sy, Inches(0.3), Inches(0.25), "✓", size=14, color=color, bold=True)
    txt(sl, sx + Inches(0.35), sy, Inches(5), Inches(0.25),
        title, size=13, color=NAVY, bold=True)
    txt(sl, sx + Inches(0.35), sy + Inches(0.3), Inches(5), Inches(0.3),
        desc, size=11, color=DARK_GRAY)

page_num(sl, 20)


# ─── Save ───
output_path = "/home/user/AI-ROI-simulator_v1/AI-ROI-Simulator_Guide.pptx"
prs.save(output_path)
print(f"Saved: {output_path} ({TOTAL_SLIDES} slides)")
